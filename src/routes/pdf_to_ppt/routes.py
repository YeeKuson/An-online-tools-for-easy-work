import os
import uuid
import zipfile
from flask import request, jsonify, current_app
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches
import io
import sys
import platform
from pdf2image import convert_from_path

from src.routes.pdf_to_ppt import pdf_to_ppt_bp

# 配置上传和下载路径
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'static', 'uploads')
DOWNLOAD_FOLDER = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'static', 'downloads')

# 确保目录存在
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

# 允许的文件扩展名
ALLOWED_EXTENSIONS = {'pdf'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def check_poppler_installed():
    """检查poppler是否已安装并在PATH中"""
    system = platform.system()
    try:
        if system == 'Windows':
            # 检查Windows下是否有poppler的bin目录在PATH中
            paths = os.environ.get('PATH', '').split(';')
            for path in paths:
                if os.path.exists(path) and 'poppler' in path.lower() and os.path.exists(os.path.join(path, 'pdftoppm.exe')):
                    return True
            return False
        else:
            # 在Linux/Mac上检查pdftoppm命令是否可用
            import subprocess
            subprocess.run(['pdftoppm', '-v'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            return True
    except (FileNotFoundError, subprocess.SubprocessError):
        return False

@pdf_to_ppt_bp.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    """
    将PDF文件转换为PPT
    
    请求参数:
    - file: PDF文件
    - quality: 图片质量 (low, medium, high, 默认为medium)
    
    返回:
    - 成功: {"success": true, "download_url": "下载链接", "page_count": 页数}
    - 失败: {"success": false, "error": "错误信息"}
    """
    # 检查是否有文件上传
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': '没有上传文件'}), 400
    
    file = request.files['file']
    
    # 检查文件名是否为空
    if file.filename == '':
        return jsonify({'success': False, 'error': '未选择文件'}), 400
    
    # 检查文件类型
    if not allowed_file(file.filename):
        return jsonify({'success': False, 'error': '不支持的文件类型，仅支持PDF文件'}), 400
    
    # 获取参数
    quality = request.form.get('quality', 'medium').lower()
    if quality not in ['low', 'medium', 'high']:
        quality = 'medium'
    
    # 设置DPI基于质量
    dpi_map = {
        'low': 100,
        'medium': 200,
        'high': 300
    }
    dpi = dpi_map[quality]
    
    # 生成唯一标识符
    unique_id = str(uuid.uuid4())
    
    try:
        # 生成安全的文件名并保存上传的文件
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, f"{unique_id}_{filename}")
        file.save(pdf_path)
        
        # 创建一个临时目录来存储图片
        images_dir = os.path.join(DOWNLOAD_FOLDER, f"{unique_id}_images")
        os.makedirs(images_dir, exist_ok=True)
        
        # 检查poppler是否已安装
        if not check_poppler_installed():
            return jsonify({'success': False, 'error': 'PDF转PPT需要安装Poppler。请访问 https://github.com/oschwartz10612/poppler-windows/releases 下载并安装，然后将bin目录添加到系统PATH环境变量中。'}), 500
        
        # 将PDF转换为图片
        images = convert_from_path(pdf_path, dpi=dpi)
        
        # 保存图片
        image_paths = []
        for i, image in enumerate(images):
            img_filename = f"page_{i+1}.jpg"
            img_path = os.path.join(images_dir, img_filename)
            image.save(img_path, 'JPEG')
            image_paths.append(img_path)
        
        page_count = len(images)
        
        # 创建PPT
        prs = Presentation()
        
        # 设置幻灯片尺寸为16:9
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # 为每个图片创建一个幻灯片
        for img_path in image_paths:
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加图片到幻灯片，并调整大小以适应幻灯片
            left = top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height
            
            slide.shapes.add_picture(img_path, left, top, width, height)
        
        # 保存PPT
        ppt_filename = f"{os.path.splitext(filename)[0]}.pptx"
        ppt_path = os.path.join(DOWNLOAD_FOLDER, f"{unique_id}_{ppt_filename}")
        prs.save(ppt_path)
        
        # 生成下载URL
        download_url = f"/download/{unique_id}_{ppt_filename}"
        
        # 删除临时文件（可选）
        # os.remove(pdf_path)
        # for img_path in image_paths:
        #     os.remove(img_path)
        # os.rmdir(images_dir)
        
        return jsonify({
            'success': True, 
            'download_url': download_url,
            'filename': ppt_filename,
            'page_count': page_count
        })
        
    except Exception as e:
        error_message = str(e)
        if "poppler" in error_message.lower():
            error_message = "PDF转PPT需要安装Poppler。请访问 https://github.com/oschwartz10612/poppler-windows/releases 下载并安装，然后将bin目录添加到系统PATH环境变量中。"
        
        return jsonify({'success': False, 'error': f'转换过程中出错: {error_message}'}), 500
